#!/usr/bin/env python3
"""Rasterize a .pptx to PNGs with Pillow so slide layout can be eyeballed.

This is a layout-QA approximation, not a faithful renderer: it draws solid
fills, outlines, pictures and greedy-wrapped text. Liberation Sans is metric
compatible with Arial, so text extents are close to what PowerPoint produces —
close enough to catch overflow, overlap and off-slide shapes.

Usage: preview_pptx.py <deck.pptx> <out_dir> [dpi]
"""

import sys
from io import BytesIO
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

EMU_PER_INCH = 914400

FONT_REGULAR = "/usr/share/fonts/liberation/LiberationSans-Regular.ttf"
FONT_BOLD = "/usr/share/fonts/liberation/LiberationSans-Bold.ttf"

_font_cache: dict[tuple[str, int], ImageFont.FreeTypeFont] = {}


def font_for(size_pt: float, bold: bool, dpi: int) -> ImageFont.FreeTypeFont:
    path = FONT_BOLD if bold else FONT_REGULAR
    px = max(1, round(size_pt * dpi / 72))
    key = (path, px)
    if key not in _font_cache:
        _font_cache[key] = ImageFont.truetype(path, px)
    return _font_cache[key]


def solid_fill(shape):
    try:
        if shape.fill.type is not None and shape.fill.type == 1:
            return tuple(shape.fill.fore_color.rgb)
    except Exception:
        pass
    return None


def outline(shape):
    try:
        ln = shape.line
        if ln.fill.type == 1:
            width = ln.width.pt if ln.width else 1.0
            return tuple(ln.color.rgb), max(1, round(width))
    except Exception:
        pass
    return None


def polygon_for(shape, box):
    """Approximate an autoshape's outline as a polygon, or None for a rect."""
    x0, y0, x1, y1 = box
    w, h = x1 - x0, y1 - y0
    try:
        kind = shape.auto_shape_type
    except Exception:
        return None

    if kind == MSO_SHAPE.DIAMOND:
        return [(x0 + w / 2, y0), (x1, y0 + h / 2), (x0 + w / 2, y1), (x0, y0 + h / 2)]
    if kind == MSO_SHAPE.CHEVRON:
        notch = min(w * 0.25, h / 2)
        return [
            (x0, y0), (x1 - notch, y0), (x1, y0 + h / 2),
            (x1 - notch, y1), (x0, y1), (x0 + notch, y0 + h / 2),
        ]
    if kind == MSO_SHAPE.PENTAGON:
        notch = min(w * 0.25, h / 2)
        return [(x0, y0), (x1 - notch, y0), (x1, y0 + h / 2), (x1 - notch, y1), (x0, y1)]
    if kind == MSO_SHAPE.RIGHT_TRIANGLE:
        return [(x0, y1), (x1, y1), (x0, y0)]
    if kind == MSO_SHAPE.ISOSCELES_TRIANGLE:
        return [(x0 + w / 2, y0), (x1, y1), (x0, y1)]
    return None


def draw_shape_body(draw, shape, box):
    fill = solid_fill(shape)
    line = outline(shape)
    if fill is None and line is None:
        return

    poly = polygon_for(shape, box)
    if poly is not None:
        draw.polygon(poly, fill=fill, outline=line[0] if line else None)
        return

    rounded = False
    try:
        rounded = shape.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE
    except Exception:
        pass
    is_oval = False
    try:
        is_oval = shape.auto_shape_type == MSO_SHAPE.OVAL
    except Exception:
        pass

    if is_oval:
        draw.ellipse(box, fill=fill, outline=line[0] if line else None,
                     width=line[1] if line else 1)
    elif rounded:
        radius = max(2, int(min(box[2] - box[0], box[3] - box[1]) * 0.12))
        draw.rounded_rectangle(box, radius=radius, fill=fill,
                              outline=line[0] if line else None,
                              width=line[1] if line else 1)
    else:
        draw.rectangle(box, fill=fill, outline=line[0] if line else None,
                       width=line[1] if line else 1)


def wrap(text, font, max_px, draw):
    """Greedy word wrap; returns list of lines."""
    if not text:
        return [""]
    lines = []
    for hard_line in text.split("\n"):
        words = hard_line.split(" ")
        current = ""
        for word in words:
            probe = word if not current else f"{current} {word}"
            if draw.textlength(probe, font=font) <= max_px or not current:
                current = probe
            else:
                lines.append(current)
                current = word
        lines.append(current)
    return lines


def draw_text_frame(draw, shape, box, scale, dpi, warnings, slide_no):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    if not tf.text.strip():
        return

    ml = (tf.margin_left or 0) * scale
    mr = (tf.margin_right or 0) * scale
    mt = (tf.margin_top or 0) * scale
    x0 = box[0] + ml
    text_width = max(8, (box[2] - box[0]) - ml - mr)
    y = box[1] + mt

    for para in tf.paragraphs:
        runs = [r for r in para.runs if r.text]
        text = "".join(r.text for r in runs)
        first = runs[0] if runs else None
        size_pt = 18.0
        bold = False
        color = (0, 0, 0)
        if first is not None:
            if first.font.size:
                size_pt = first.font.size.pt
            bold = bool(first.font.bold)
            try:
                if first.font.color and first.font.color.rgb:
                    color = tuple(first.font.color.rgb)
            except Exception:
                pass

        font = font_for(size_pt, bold, dpi)
        line_h = size_pt * 1.21 * dpi / 72
        y += (para.space_before.pt if para.space_before else 0) * dpi / 72

        for line in wrap(text, font, text_width, draw):
            lx = x0
            if para.alignment == PP_ALIGN.CENTER:
                lx = x0 + (text_width - draw.textlength(line, font=font)) / 2
            elif para.alignment == PP_ALIGN.RIGHT:
                lx = x0 + text_width - draw.textlength(line, font=font)
            draw.text((lx, y), line, font=font, fill=color)
            y += line_h

        y += (para.space_after.pt if para.space_after else 0) * dpi / 72

    overflow = y - box[3]
    if overflow > line_h * 0.6:
        warnings.append(
            f"slide {slide_no}: text overflows its box by ~{overflow / dpi:.2f}in "
            f"-> {tf.text.strip()[:60]!r}"
        )


def draw_connector(draw, shape, scale, ox, oy):
    xfrm = shape._element.spPr.find(qn("a:xfrm"))
    flip_h = xfrm is not None and xfrm.get("flipH") == "1"
    flip_v = xfrm is not None and xfrm.get("flipV") == "1"
    x0 = ox + shape.left * scale
    y0 = oy + shape.top * scale
    x1 = x0 + shape.width * scale
    y1 = y0 + shape.height * scale
    start = (x1 if flip_h else x0, y1 if flip_v else y0)
    end = (x0 if flip_h else x1, y0 if flip_v else y1)
    line = outline(shape)
    color = line[0] if line else (0, 0, 0)
    width = line[1] if line else 1
    draw.line([start, end], fill=color, width=width)
    # Arrow head marker so direction is visible in the preview.
    draw.ellipse(
        [end[0] - width * 2, end[1] - width * 2, end[0] + width * 2, end[1] + width * 2],
        fill=color,
    )


def render(path, out_dir, dpi=120):
    prs = Presentation(path)
    scale = dpi / EMU_PER_INCH
    width_px = round(prs.slide_width * scale)
    height_px = round(prs.slide_height * scale)
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    warnings = []
    written = []
    for index, slide in enumerate(prs.slides, start=1):
        img = Image.new("RGB", (width_px, height_px), "white")
        draw = ImageDraw.Draw(img)

        for shape in slide.shapes:
            if shape.left is None or shape.top is None:
                continue
            box = (
                shape.left * scale,
                shape.top * scale,
                (shape.left + shape.width) * scale,
                (shape.top + shape.height) * scale,
            )
            if shape.shape_type == MSO_SHAPE_TYPE.LINE:
                draw_connector(draw, shape, scale, 0, 0)
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    pic = Image.open(BytesIO(shape.image.blob)).convert("RGB")
                    size = (max(1, round(box[2] - box[0])), max(1, round(box[3] - box[1])))
                    img.paste(pic.resize(size), (round(box[0]), round(box[1])))
                except Exception:
                    draw.rectangle(box, outline=(200, 200, 200))
                continue

            draw_shape_body(draw, shape, box)
            draw_text_frame(draw, shape, box, scale, dpi, warnings, index)

            if box[0] < -1 or box[1] < -1 or box[2] > width_px + 1 or box[3] > height_px + 1:
                warnings.append(f"slide {index}: shape extends past the slide edge {box}")

        target = out_dir / f"slide{index}.png"
        img.save(target)
        written.append(str(target))

    for line in warnings:
        print("WARN:", line)
    print(f"rendered {len(written)} slides -> {out_dir}")
    return written


if __name__ == "__main__":
    render(sys.argv[1], sys.argv[2], int(sys.argv[3]) if len(sys.argv) > 3 else 120)
