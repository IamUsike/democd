#!/usr/bin/env python3
"""Build the six-slide client deck for the Transaction Monitoring & Alerts project.

Visual language follows the HSBC-style template in runtime_rebels.pptx.pdf:
deep-red brand panel bookends, generous whitespace, hairline-bordered cards.
Every figure quoted here is traceable to the repo or to docs/load-test-results.md.

Run: .pptx-venv/bin/python scripts/build_final_pptx.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ── Palette ───────────────────────────────────────────────────────────────
# RED_DEEP matches the background baked into hsbc.png so the logo sits
# invisibly on the brand panel.
RED_DEEP = RGBColor(0x73, 0x00, 0x14)
RED = RGBColor(0xC4, 0x1E, 0x3A)
RED_MID = RGBColor(0x9B, 0x0F, 0x28)
INK = RGBColor(0x14, 0x18, 0x1D)
BODY = RGBColor(0x39, 0x41, 0x4C)
MUTED = RGBColor(0x7C, 0x86, 0x95)
HAIRLINE = RGBColor(0xD9, 0xDE, 0xE4)
MIST = RGBColor(0xF4, 0xF6, 0xF8)
CARD = RGBColor(0xEA, 0xEE, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x1B, 0x7A, 0x3D)
AMBER = RGBColor(0xB0, 0x59, 0x0A)

TEAM = "AGILE-ish"
MEMBERS = "Shreya  ·  Sathwik  ·  Rameez"
PROJECT = "Transaction Monitoring & Alerts"
TOTAL_SLIDES = 6

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = 0.6
CONTENT_W = 12.133

REPO = Path(__file__).resolve().parent.parent
ASSETS = REPO / "scripts" / "assets"
OUT = str(REPO / "FINAL-PRESENTATION.pptx")


def asset(name):
    path = ASSETS / name
    return str(path) if path.exists() else None


# ── Primitives ────────────────────────────────────────────────────────────
def style_run(run, size, bold=False, color=INK, spacing=None, font="Arial"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    if spacing:
        run.font._rPr.set("spc", str(int(spacing * 100)))


def text(slide, x, y, w, h, body, size=12, bold=False, color=BODY,
         align=PP_ALIGN.LEFT, spacing=None, line_spacing=None):
    """Add a word-wrapped textbox. `body` may contain newlines."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    for index, line in enumerate(body.split("\n")):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.alignment = align
        para.space_before = Pt(0)
        para.space_after = Pt(0)
        if line_spacing:
            para.line_spacing = line_spacing
        run = para.add_run()
        run.text = line
        style_run(run, size, bold, color, spacing)
    return box


def shape(slide, kind, x, y, w, h, fill=None, line=None, line_w=0.75, dash=None):
    item = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        item.fill.background()
    else:
        item.fill.solid()
        item.fill.fore_color.rgb = fill
    if line is None:
        item.line.fill.background()
    else:
        item.line.color.rgb = line
        item.line.width = Pt(line_w)
        if dash:
            item.line.dash_style = dash
    item.shadow.inherit = False
    return item


def rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75, dash=None):
    return shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, fill, line, line_w, dash)


def panel(slide, x, y, w, h, fill=MIST, line=HAIRLINE, radius=0.06):
    """A soft card. Rounded corners are set via the shape's adjustment value."""
    item = shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, line)
    try:
        item.adjustments[0] = radius / min(w, h) if min(w, h) else 0.1
    except Exception:
        pass
    return item


def hairline(slide, x, y, w, color=HAIRLINE):
    return rect(slide, x, y, w, 0.012, fill=color)


def accent_rule(slide, x, y, w=1.15, color=RED, thickness=0.038):
    return rect(slide, x, y, w, thickness, fill=color)


def arrow(slide, x1, y1, x2, y2, color=RED, width=1.25, dash=None):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    if dash:
        line.line.dash_style = dash
    ln = line._element.spPr.get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    return line


def centered(item, lines):
    """Fill an autoshape with vertically centred lines: (text, size, bold, color)."""
    frame = item.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    for index, (body, size, bold, color) in enumerate(lines):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        para.space_before = Pt(2 if index else 0)
        para.space_after = Pt(0)
        run = para.add_run()
        run.text = body
        style_run(run, size, bold, color)


def dot_field(slide, x, y, cols, rows, step=0.17, size=0.075):
    """Faint diamond grid echoing the template's brand pattern."""
    for row in range(rows):
        for col in range(cols):
            if (row + col) % 2:
                continue
            dot = shape(
                slide, MSO_SHAPE.DIAMOND,
                x + col * step, y + row * step, size, size, fill=RED,
            )
            fade = 0.30 - 0.022 * (row + col)
            set_alpha(dot, max(0.06, fade))


def set_alpha(item, alpha):
    fill = item._element.spPr.find(qn("a:solidFill"))
    if fill is None:
        return
    color = fill.find(qn("a:srgbClr"))
    if color is None:
        return
    color.append(color.makeelement(qn("a:alpha"), {"val": str(int(alpha * 100000))}))


def header(slide, eyebrow, title, subtitle=None, title_size=29):
    text(slide, MARGIN, 0.44, 8.0, 0.24, eyebrow, size=9.5, bold=True,
         color=RED, spacing=1.6)
    text(slide, MARGIN, 0.70, 11.6, 0.5, title, size=title_size, bold=True, color=INK)
    accent_rule(slide, MARGIN, 1.24)
    if subtitle:
        text(slide, MARGIN, 1.40, 11.4, 0.3, subtitle, size=11.5, color=MUTED)


def footer(slide, number):
    hairline(slide, MARGIN, 6.99, CONTENT_W)
    text(slide, MARGIN, 7.08, 8.0, 0.24, f"{TEAM}  ·  {PROJECT}", size=9, color=MUTED)
    text(slide, 11.0, 7.08, 1.73, 0.24, f"{number} / {TOTAL_SLIDES}", size=9,
         color=MUTED, align=PP_ALIGN.RIGHT)


def stat_band(slide, y, cells, height=0.95):
    """Full-width row of big numbers separated by hairlines."""
    panel(slide, MARGIN, y, CONTENT_W, height, fill=MIST)
    cell_w = CONTENT_W / len(cells)
    for index, (value, label) in enumerate(cells):
        cx = MARGIN + index * cell_w
        if index:
            rect(slide, cx, y + 0.16, 0.008, height - 0.32, fill=HAIRLINE)
        text(slide, cx + 0.28, y + 0.16, cell_w - 0.5, 0.4, value, size=23,
             bold=True, color=RED_DEEP)
        text(slide, cx + 0.28, y + 0.60, cell_w - 0.4, 0.24, label, size=10,
             color=MUTED)


# ── Slide 1 — title ───────────────────────────────────────────────────────
def slide_title(prs, blank):
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, 4.95, 7.5, fill=RED_DEEP)
    dot_field(s, 3.45, 4.50, cols=8, rows=12)
    rect(s, 4.95, 0, 0.055, 7.5, fill=RED)

    text(s, 0.55, 0.62, 3.9, 0.24, "CLIENT REVIEW  ·  AUGUST 2026", size=9.5,
         bold=True, color=RGBColor(0xE8, 0xA6, 0xB2), spacing=1.6)
    text(s, 0.55, 1.05, 4.0, 2.3, "TRANSACTION\nMONITORING\n& ALERTS",
         size=31, bold=True, color=WHITE, line_spacing=1.14)
    accent_rule(s, 0.55, 2.95, w=1.5, color=WHITE, thickness=0.03)
    text(s, 0.55, 3.20, 3.85, 0.9,
         "Catch the risky payment while it is still\nin flight — then give one operator a\nclear path to close it.",
         size=12.5, color=RGBColor(0xF2, 0xD6, 0xDB), line_spacing=1.3)

    text(s, 0.55, 4.92, 3.9, 0.22, "PRESENTED BY TEAM", size=9, bold=True,
         color=RGBColor(0xD8, 0x92, 0xA0), spacing=1.4)
    text(s, 0.55, 5.16, 3.9, 0.5, TEAM, size=27, bold=True, color=WHITE)
    text(s, 0.55, 5.78, 3.9, 0.22, "TEAM MEMBERS", size=9, bold=True,
         color=RGBColor(0xD8, 0x92, 0xA0), spacing=1.4)
    text(s, 0.55, 6.02, 3.9, 0.26, MEMBERS, size=12.5, color=WHITE)

    if asset("hsbc.png"):
        s.shapes.add_picture(asset("hsbc.png"), Inches(0.55), Inches(6.68),
                             height=Inches(0.4))
    else:
        text(s, 0.55, 6.72, 2.0, 0.3, "HSBC", size=13, bold=True, color=WHITE)

    if asset("plaza.png"):
        s.shapes.add_picture(asset("plaza.png"), Inches(8.55), 0,
                             width=Inches(4.783), height=Inches(7.5))
    else:
        rect(s, 8.55, 0, 4.783, 7.5, fill=CARD)

    if asset("how_we_lead.png"):
        s.shapes.add_picture(asset("how_we_lead.png"), Inches(5.74),
                             Inches(2.05), height=Inches(1.85))

    text(s, 5.55, 4.55, 2.85, 0.22, "THE OPERATOR LOOP", size=9, bold=True,
         color=RED, spacing=1.4)
    loop = ["Ingest", "Evaluate", "Alert", "Investigate", "Close"]
    ly = 4.86
    for index, step in enumerate(loop):
        rect(s, 5.55, ly + 0.02, 0.075, 0.16, fill=RED if index < 3 else HAIRLINE)
        text(s, 5.78, ly, 2.6, 0.22, step, size=11.5,
             bold=index < 3, color=INK if index < 3 else MUTED)
        ly += 0.33

    return s


# ── Slide 2 — problem & solution ──────────────────────────────────────────
def slide_problem(prs, blank):
    s = prs.slides.add_slide(blank)
    header(
        s, "THE BRIEF",
        "Payments do not wait for a night-time batch",
        "Banks and merchants push transactions continuously. Risk has to surface while the money is moving.",
    )

    text(s, MARGIN, 1.92, 5.4, 0.26, "What makes it hard", size=15, bold=True,
         color=RED_DEEP)
    problems = [
        ("01", "Volume never pauses",
         "Two simulated feeds plus a load harness can push hundreds of transactions a second into one API."),
        ("02", "Manual triage does not scale",
         "Amount spikes are the easy case. Velocity bursts, first-time payees and creeping daily totals hide in the noise."),
        ("03", "Detection alone is not an answer",
         "An operator needs one queue, one audit trail, and a status that cannot be skipped or reversed by accident."),
    ]
    y = 2.34
    for tag, title, body in problems:
        text(s, MARGIN, y, 0.5, 0.3, tag, size=16, bold=True, color=RED)
        text(s, MARGIN + 0.62, y + 0.02, 4.75, 0.26, title, size=13, bold=True,
             color=INK)
        text(s, MARGIN + 0.62, y + 0.32, 4.75, 0.66, body, size=11, color=MUTED,
             line_spacing=1.28)
        y += 1.12

    text(s, 6.95, 1.92, 5.8, 0.26, "What we delivered", size=15, bold=True,
         color=RED_DEEP)
    text(s, 6.95, 2.26, 5.78, 0.58,
         "One Spring Boot deployable and one operator dashboard: ingest, evaluate against four rules in the same request, then work the alert to a close.",
         size=12, bold=True, color=INK, line_spacing=1.3)

    cards = [
        ("Four rules, live-configurable",
         "Thresholds are edited in the UI and stored in the database — no redeploy."),
        ("Soft multi-source tenancy",
         "One schema; every row carries source_type, source_id and source_name."),
        ("Measured, not assumed",
         "Three k6 passes, EXPLAIN on every hot query, timers on the engine itself."),
    ]
    cy = 3.06
    for title, body in cards:
        panel(s, 6.95, cy, 5.78, 0.80, fill=MIST)
        rect(s, 6.95, cy, 0.055, 0.80, fill=RED)
        text(s, 7.18, cy + 0.14, 5.4, 0.24, title, size=12, bold=True, color=INK)
        text(s, 7.18, cy + 0.42, 5.4, 0.3, body, size=10.5, color=MUTED,
             line_spacing=1.2)
        cy += 0.90

    stat_band(s, 5.86, [
        ("4", "detection rules, all runtime-configurable"),
        ("5", "alert states with validated transitions"),
        ("14", "REST endpoints under /api/v1"),
        ("89", "unit tests across 17 classes"),
    ])
    footer(s, 2)
    return s


# ── Slide 3 — architecture ────────────────────────────────────────────────
def slide_architecture(prs, blank):
    s = prs.slides.add_slide(blank)
    header(
        s, "ARCHITECTURE",
        "A modular monolith with one seam that matters",
        "One deployable, clean package boundaries — and a rule engine kept behind an interface so it can leave later.",
    )

    col_a_x, col_a_w = MARGIN, 2.15
    vm_x, vm_w = 3.62, 5.93
    col_c_x, col_c_w = 9.95, 2.78

    # Column A — clients
    text(s, col_a_x, 1.86, col_a_w, 0.22, "CLIENTS & FEEDS", size=9, bold=True,
         color=MUTED, spacing=1.3)
    sources = [
        ("Operator browser", "Loads the dashboard,\ndrives the alert lifecycle"),
        ("Feed simulators", "Go service, scenario packs\nPOST /api/v1/transactions"),
        ("k6 load harness", "Ramp, mixed and soak\nprofiles up to 200 VUs"),
    ]
    sy = 2.15
    source_centres = []
    for title, body in sources:
        panel(s, col_a_x, sy, col_a_w, 0.95, fill=WHITE, line=HAIRLINE)
        text(s, col_a_x + 0.16, sy + 0.15, col_a_w - 0.32, 0.22, title,
             size=11, bold=True, color=INK)
        text(s, col_a_x + 0.16, sy + 0.42, col_a_w - 0.28, 0.44, body,
             size=8.8, color=MUTED, line_spacing=1.2)
        source_centres.append(sy + 0.475)
        sy += 1.10

    # Column B — the application VM
    panel(s, vm_x, 1.98, vm_w, 3.42, fill=MIST)
    text(s, vm_x + 0.18, 2.10, 4.6, 0.22, "APPLICATION VM  ·  DOCKER COMPOSE",
         size=9, bold=True, color=RED_DEEP, spacing=1.3)

    nginx = panel(s, vm_x + 0.18, 2.42, vm_w - 0.36, 0.56, fill=WHITE)
    centered(nginx, [("nginx  +  React 19 SPA   ·   :8082", 11, True, INK)])

    api_y = 3.22
    panel(s, vm_x + 0.18, api_y, vm_w - 0.36, 1.98, fill=WHITE)
    text(s, vm_x + 0.36, api_y + 0.12, 4.4, 0.22,
         "SPRING BOOT 3.4.4  ·  JAVA 21  ·  :8081", size=9, bold=True,
         color=RED_DEEP, spacing=1.3)

    modules = [
        ("api", "controllers + DTOs", False),
        ("transaction", "record, query, filters", False),
        ("rule", "RuleEngine + 4 rules", True),
        ("alert", "lifecycle + transitions", False),
    ]
    chip_w = (vm_w - 0.36 - 0.36 - 0.22) / 2
    for index, (name, body, highlight) in enumerate(modules):
        cx = vm_x + 0.36 + (index % 2) * (chip_w + 0.22)
        cy = api_y + 0.42 + (index // 2) * 0.72
        panel(s, cx, cy, chip_w, 0.62,
              fill=WHITE if highlight else MIST,
              line=RED if highlight else HAIRLINE)
        text(s, cx + 0.14, cy + 0.11, chip_w - 0.28, 0.22, name, size=10.5,
             bold=True, color=RED if highlight else INK)
        text(s, cx + 0.14, cy + 0.34, chip_w - 0.24, 0.22, body, size=8.8,
             color=MUTED)

    text(s, vm_x + 0.36, api_y + 1.80, chip_w + 1.4, 0.2,
         "extractable later — same interface, new deployable", size=8,
         bold=True, color=RED)

    # Column C — data
    text(s, col_c_x, 1.86, col_c_w, 0.22, "DATA", size=9, bold=True,
         color=MUTED, spacing=1.3)
    panel(s, col_c_x, 2.15, col_c_w, 1.95, fill=WHITE)
    text(s, col_c_x + 0.16, 2.30, col_c_w - 0.32, 0.22, "MySQL 8  ·  :3306",
         size=11, bold=True, color=INK)
    text(s, col_c_x + 0.16, 2.56, col_c_w - 0.32, 0.2,
         "reached over JPA + HikariCP", size=8.8, color=MUTED)
    tables = ["transactions", "alerts", "alert_transactions", "rule_configs"]
    ty = 2.86
    for table in tables:
        rect(s, col_c_x + 0.16, ty + 0.055, 0.06, 0.06, fill=RED)
        text(s, col_c_x + 0.32, ty, col_c_w - 0.48, 0.2, table, size=9.5,
             color=BODY)
        ty += 0.27

    panel(s, col_c_x, 4.25, col_c_w, 1.15, fill=MIST)
    text(s, col_c_x + 0.16, 4.40, col_c_w - 0.32, 0.22, "7 named indexes",
         size=11, bold=True, color=INK)
    text(s, col_c_x + 0.16, 4.66, col_c_w - 0.3, 0.6,
         "Account+time, account+payee and\nsource+time paths confirmed with\nEXPLAIN before blaming the schema.",
         size=8.8, color=MUTED, line_spacing=1.18)

    # Wiring
    arrow(s, col_a_x + col_a_w, source_centres[0], vm_x + 0.14, 2.70)
    arrow(s, col_a_x + col_a_w, source_centres[1], vm_x + 0.14, source_centres[1])
    arrow(s, col_a_x + col_a_w, source_centres[2], vm_x + 0.14, source_centres[2])
    arrow(s, vm_x + 0.9, 2.98, vm_x + 0.9, api_y - 0.02, color=MUTED, width=1.0)
    text(s, vm_x + 1.02, 3.02, 1.6, 0.2, "REST /api/v1", size=8, color=MUTED)
    arrow(s, vm_x + vm_w - 0.14, 3.6, col_c_x + 0.04, 3.6)

    # Synchronous request path
    text(s, MARGIN, 5.62, 7.0, 0.22, "SYNCHRONOUS DETECTION PATH", size=9,
         bold=True, color=RED, spacing=1.4)
    steps = [
        "POST /transactions",
        "Persist transaction",
        "RuleEngine.evaluate()",
        "Create OPEN alert(s)",
        "201 Created",
    ]
    pill_w = (CONTENT_W - 4 * 0.34) / 5
    for index, label in enumerate(steps):
        px = MARGIN + index * (pill_w + 0.34)
        pill = panel(s, px, 5.90, pill_w, 0.52,
                     fill=RED_DEEP if index == 2 else MIST,
                     line=RED_DEEP if index == 2 else HAIRLINE)
        centered(pill, [(f"{index + 1}   {label}", 10.5, True,
                         WHITE if index == 2 else INK)])
        if index < 4:
            arrow(s, px + pill_w + 0.06, 6.16, px + pill_w + 0.30, 6.16,
                  color=MUTED, width=1.0)

    text(s, MARGIN, 6.56, CONTENT_W, 0.34,
         "Detection finishes before the API answers, so time-to-detect is bounded by the request itself. "
         "No queue in this path yet — that is the deliberate next change, not a rewrite.",
         size=10.5, color=MUTED, line_spacing=1.2)
    footer(s, 3)
    return s


# ── Slide 4 — detection & lifecycle ───────────────────────────────────────
def slide_detection(prs, blank):
    s = prs.slides.add_slide(blank)
    header(
        s, "DETECTION & LIFECYCLE",
        "Four rules, five states, one audit trail",
        "Adding a rule type means adding one class. The engine and the Rule interface never change.",
    )

    text(s, MARGIN, 1.92, 5.9, 0.26, "The rule set", size=15, bold=True,
         color=RED_DEEP)
    rules = [
        ("Amount Threshold", "Single transaction above a ceiling",
         "> 10,000", "MEDIUM"),
        ("Velocity", "Too many transactions in a rolling window",
         "> 5 in 10 min", "LOW"),
        ("New Payee", "First time this account pays this payee",
         "0 prior rows", "LOW"),
        ("Daily Limit", "Cumulative debits for the calendar day",
         "> 50,000", "MEDIUM"),
    ]
    card_w = (5.95 - 0.25) / 2
    for index, (name, what, default, severity) in enumerate(rules):
        cx = MARGIN + (index % 2) * (card_w + 0.25)
        cy = 2.32 + (index // 2) * 1.30
        panel(s, cx, cy, card_w, 1.16, fill=MIST)
        text(s, cx + 0.18, cy + 0.14, card_w - 0.36, 0.24, name, size=12,
             bold=True, color=INK)
        text(s, cx + 0.18, cy + 0.41, card_w - 0.32, 0.42, what, size=9.5,
             color=MUTED, line_spacing=1.2)
        text(s, cx + 0.18, cy + 0.85, 1.6, 0.22, default, size=10.5, bold=True,
             color=RED)
        text(s, cx + card_w - 1.05, cy + 0.85, 0.87, 0.22, severity, size=8.5,
             bold=True, color=MUTED, align=PP_ALIGN.RIGHT)

    text(s, MARGIN, 5.02, 5.9, 0.44,
         "Every default lives in rule_configs and is editable from the Rules screen — "
         "no redeploy to retune a threshold.",
         size=10.5, color=MUTED, line_spacing=1.22)

    # Lifecycle
    text(s, 6.95, 1.92, 5.8, 0.26, "Alert lifecycle", size=15, bold=True,
         color=RED_DEEP)
    states = [("OPEN", RED_DEEP), ("ACKNOWLEDGED", RED_MID),
              ("INVESTIGATING", RED), ("CLOSED", GREEN)]
    state_w = (5.78 - 3 * 0.20) / 4
    centres = []
    for index, (name, color) in enumerate(states):
        sx = 6.95 + index * (state_w + 0.20)
        pill = panel(s, sx, 2.34, state_w, 0.5, fill=color, line=color)
        centered(pill, [(name, 8.5, True, WHITE)])
        centres.append(sx + state_w / 2)
        if index < 3:
            arrow(s, sx + state_w + 0.02, 2.59, sx + state_w + 0.17, 2.59,
                  color=MUTED, width=1.0)

    dismissed = panel(s, centres[1] + 0.10, 3.34, 1.9, 0.5, fill=WHITE,
                      line=AMBER)
    centered(dismissed, [("DISMISSED", 8.5, True, AMBER)])
    arrow(s, centres[1], 2.86, centres[1] + 0.35, 3.32, color=AMBER, width=1.0)
    arrow(s, centres[2], 2.86, centres[1] + 1.72, 3.32, color=AMBER, width=1.0)

    text(s, 6.95, 3.98, 5.78, 0.44,
         "CLOSED and DISMISSED are terminal. Anything else raises InvalidAlertTransitionException — "
         "checked in the service, not the controller.",
         size=10, color=MUTED, line_spacing=1.22)

    panel(s, 6.95, 4.56, 5.78, 0.9, fill=MIST)
    rect(s, 6.95, 4.56, 0.055, 0.9, fill=RED)
    text(s, 7.18, 4.70, 5.4, 0.22, "When two rules fire on one transaction",
         size=11, bold=True, color=INK)
    text(s, 7.18, 4.96, 5.4, 0.4,
         "One alert, both reasons listed, severity escalated to HIGH.",
         size=10, color=MUTED)

    # Quality bar
    panel(s, MARGIN, 5.72, CONTENT_W, 1.06, fill=WHITE, line=HAIRLINE)
    text(s, MARGIN + 0.28, 5.88, 6.0, 0.22, "How we kept it honest", size=11.5,
         bold=True, color=RED_DEEP)
    proofs = [
        ("Test first", "Failing test before every rule and every illegal transition"),
        ("Coverage", "89 unit tests across 17 classes, H2 for repository tests"),
        ("Contract", "OpenAPI published at /swagger-ui.html"),
        ("Pipeline", "Jenkins rebuilds the whole Compose stack on push to dev"),
    ]
    proof_w = (CONTENT_W - 0.56) / 4
    for index, (label, body) in enumerate(proofs):
        px = MARGIN + 0.28 + index * proof_w
        text(s, px, 6.20, proof_w - 0.2, 0.2, label, size=9, bold=True,
             color=RED, spacing=1.2)
        text(s, px, 6.42, proof_w - 0.22, 0.34, body, size=9, color=MUTED,
             line_spacing=1.18)

    footer(s, 4)
    return s


# ── Slide 5 — performance & the blocker ───────────────────────────────────
def slide_performance(prs, blank):
    s = prs.slides.add_slide(blank)
    header(
        s, "EVIDENCE",
        "We pushed it until it hurt, then found out why",
        "k6 from a Windows client against the Linux app VM — JVM and MySQL sharing roughly 2 vCPU and 3.7 GiB.",
    )

    metrics = [
        ("PASS 1  ·  WRITE RAMP", "234", "requests/sec",
         "200 VUs  ·  p95 763 ms  ·  0% failed", GREEN),
        ("PASS 2  ·  MIXED 80/20", "242", "requests/sec",
         "150 VUs  ·  p95 623 ms  ·  0% failed", GREEN),
        ("PASS 3  ·  10-MINUTE SOAK", "214", "requests/sec",
         "140 VUs  ·  p95 1.13 s  ·  0.09% failed", AMBER),
    ]
    card_w = (CONTENT_W - 2 * 0.26) / 3
    for index, (label, value, unit, detail, accent) in enumerate(metrics):
        cx = MARGIN + index * (card_w + 0.26)
        panel(s, cx, 1.86, card_w, 1.42, fill=MIST)
        rect(s, cx, 1.86, card_w, 0.05, fill=accent)
        text(s, cx + 0.24, 2.04, card_w - 0.4, 0.2, label, size=8.5, bold=True,
             color=MUTED, spacing=1.2)
        text(s, cx + 0.24, 2.28, 1.5, 0.5, value, size=30, bold=True, color=INK)
        text(s, cx + 1.35, 2.50, 1.6, 0.24, unit, size=10.5, color=MUTED)
        text(s, cx + 0.24, 2.90, card_w - 0.4, 0.24, detail, size=9.5,
             color=BODY)

    base_y = 5.62
    max_h = 1.62

    text(s, MARGIN, 3.52, 4.4, 0.22, "p95 LATENCY BY PASS  (ms)", size=9,
         bold=True, color=RED, spacing=1.3)
    for index, (label, value) in enumerate([("Pass 1", 763), ("Pass 2", 623),
                                            ("Pass 3", 1130)]):
        h = max_h * value / 1200
        bx = MARGIN + 0.1 + index * 1.42
        rect(s, bx, base_y - h, 0.95, h, fill=AMBER if value > 1000 else RED_DEEP)
        text(s, bx - 0.2, base_y - h - 0.26, 1.35, 0.22, f"{value:,}", size=10,
             bold=True, color=INK, align=PP_ALIGN.CENTER)
        text(s, bx - 0.2, base_y + 0.07, 1.35, 0.2, label, size=9, color=MUTED,
             align=PP_ALIGN.CENTER)
    hairline(s, MARGIN, base_y, 4.35)

    text(s, 5.25, 3.52, 4.1, 0.22, "rule.evaluate MEAN  (ms)", size=9,
         bold=True, color=RED, spacing=1.3)
    evals = [("after P1", 57), ("after P2", 53), ("mid soak", 91),
             ("end soak", 126)]
    for index, (label, value) in enumerate(evals):
        h = max_h * value / 140
        bx = 5.35 + index * 1.02
        rect(s, bx, base_y - h, 0.68, h, fill=RED if index < 2 else AMBER)
        text(s, bx - 0.16, base_y - h - 0.26, 1.0, 0.22, str(value), size=10,
             bold=True, color=INK, align=PP_ALIGN.CENTER)
        text(s, bx - 0.16, base_y + 0.07, 1.0, 0.2, label, size=8.5,
             color=MUTED, align=PP_ALIGN.CENTER)
    hairline(s, 5.25, base_y, 4.05)

    text(s, MARGIN, 6.04, 8.7, 0.62,
         "On the short passes p95 sat at 600–800 ms while rule evaluation averaged 53–57 ms — "
         "so the latency was never in the rule engine. Under soak both curves climbed together, "
         "which is what CPU contention looks like, not a missing index.",
         size=10.5, color=BODY, line_spacing=1.24)

    # The blocker
    panel(s, 9.55, 3.44, 3.18, 3.28, fill=MIST)
    rect(s, 9.55, 3.44, 3.18, 0.05, fill=RED_DEEP)
    text(s, 9.77, 3.62, 2.8, 0.22, "THE BLOCKER", size=9, bold=True,
         color=RED_DEEP, spacing=1.4)
    text(s, 9.77, 3.86, 2.8, 0.44, "Moving MySQL to its own VM", size=12.5,
         bold=True, color=INK, line_spacing=1.16)
    text(s, 9.77, 4.34, 2.82, 0.34,
         "The soak said co-location was the cost. The fix would not connect:",
         size=9.5, color=MUTED, line_spacing=1.2)

    checks = [
        ("Windows → Linux", True),
        ("Linux A → Linux B", False),
        ("nc listener ← Windows", True),
        ("nc listener ← Linux A", False),
    ]
    cy = 4.86
    for label, ok in checks:
        text(s, 9.77, cy, 2.3, 0.2, label, size=9, color=BODY)
        text(s, 11.95, cy, 0.6, 0.2, "works" if ok else "fails", size=9,
             bold=True, color=GREEN if ok else RED, align=PP_ALIGN.RIGHT)
        cy += 0.26
    hairline(s, 9.77, cy + 0.04, 2.78)
    text(s, 9.77, cy + 0.18, 2.82, 0.5,
         "Plain nc fails the same way, so this is a security-group or routing rule — not Spring, not Docker.",
         size=9, color=MUTED, line_spacing=1.2)

    footer(s, 4)
    return s


# ── Slide 5 — future scopes ───────────────────────────────────────────────
def slide_future(prs, blank):
    s = prs.slides.add_slide(blank)
    header(
        s, "FUTURE SCOPES",
        "The seams are in place — here is what goes through them",
        "Every item below is a deployment change, not a rewrite. The interfaces on the architecture slide were designed for this.",
    )

    scopes = [
        ("01", "MySQL on its own VM",
         "Clear the Linux-to-Linux security-group rule the soak exposed, move the database off the app box, and re-run the 140 VU soak to publish the before-and-after."),
        ("02", "Async queue between ingest and evaluation",
         "Accept the transaction, answer 201, and evaluate from a queue — detection leaves the request path, so p95 stays flat as volume grows."),
        ("03", "Rule engine as its own deployable",
         "The rule package already sits behind one interface. Extract it, run N workers against the queue, and scale detection independently of the API."),
        ("04", "API keys for ingest validation",
         "Issue a key per bank and merchant feed so /api/v1/transactions authenticates its callers — soft tenancy gets a hard edge."),
        ("05", "Hardening for real money",
         "TLS end to end, masking of account and payee fields in the UI, and credentials out of source into a secrets manager."),
    ]
    y = 1.90
    for tag, title, body in scopes:
        text(s, MARGIN, y, 0.5, 0.3, tag, size=15, bold=True, color=RED)
        text(s, MARGIN + 0.60, y + 0.02, 6.9, 0.24, title, size=12.5,
             bold=True, color=INK)
        text(s, MARGIN + 0.60, y + 0.30, 6.9, 0.56, body, size=10.5,
             color=MUTED, line_spacing=1.24)
        y += 0.99

    # Target architecture — the seam carried over from the architecture slide.
    panel_x, panel_w = 8.35, 4.38
    panel(s, panel_x, 1.90, panel_w, 4.72, fill=MIST)
    text(s, panel_x + 0.22, 2.04, panel_w - 0.44, 0.22,
         "TARGET SHAPE  ·  FROM THE ARCHITECTURE SLIDE", size=8.5, bold=True,
         color=RED_DEEP, spacing=1.2)

    box_x = panel_x + 0.25
    box_w = panel_w - 0.50
    mid_x = panel_x + panel_w / 2
    stages = [
        ("Spring Boot API  ·  :8081", "accepts, persists, answers 201",
         WHITE, HAIRLINE, INK, MUTED),
        ("MESSAGE QUEUE", "ingest decoupled from detection",
         RED_DEEP, RED_DEEP, WHITE, RGBColor(0xF2, 0xD6, 0xDB)),
        ("Rule workers  ×  N", "extracted RuleEngine, scaled horizontally",
         WHITE, RED, RED, MUTED),
        ("MySQL 8  ·  own VM", "co-location contention gone",
         WHITE, HAIRLINE, INK, MUTED),
    ]
    by = 2.36
    for index, (title, body, fill, line, tcol, bcol) in enumerate(stages):
        pill = panel(s, box_x, by, box_w, 0.64, fill=fill, line=line)
        centered(pill, [(title, 10.5, True, tcol), (body, 8.5, False, bcol)])
        if index < 3:
            arrow(s, mid_x, by + 0.66, mid_x, by + 0.86, color=MUTED, width=1.0)
        by += 0.88

    text(s, box_x, by + 0.02, box_w, 0.72,
         "Same Rule interface as today — extraction is the seam highlighted "
         "on the architecture slide, and the monolith keeps an internal "
         "alert-creation endpoint.",
         size=9, color=MUTED, line_spacing=1.22)

    footer(s, 5)
    return s


# ── Slide 6 — how we worked + close ───────────────────────────────────────
def slide_close(prs, blank):
    s = prs.slides.add_slide(blank)
    rect(s, 8.40, 0, 4.933, 7.5, fill=RED_DEEP)
    dot_field(s, 11.35, 5.30, cols=8, rows=11)
    rect(s, 8.40, 0, 0.055, 7.5, fill=RED)

    text(s, MARGIN, 0.44, 7.0, 0.24, "HOW WE WORKED", size=9.5, bold=True,
         color=RED, spacing=1.6)
    text(s, MARGIN, 0.70, 7.4, 0.5, "Agile, with receipts", size=29, bold=True,
         color=INK)
    accent_rule(s, MARGIN, 1.24)

    practices = [
        ("MVP before anything clever",
         "Transactions, then the synchronous Amount rule, then the full alert lifecycle, then the UI. The other three rules waited until that path worked end to end."),
        ("Two sprints, daily stand-ups",
         "29 July to 11 August 2026. The stand-up log and three retrospectives are checked into docs/ alongside the code."),
        ("Kanban with a WIP limit of one",
         "A card moves to Review when the PR opens and to Done on merge plus a milestone tick — so status is never a guess."),
        ("Plans written before code",
         "Thirteen short design notes in .cursor/plans/, one per non-trivial feature."),
        ("Test-first on the risky parts",
         "A failing test before every rule and every illegal lifecycle transition — 89 unit tests across 17 classes, contract published at /swagger-ui.html."),
    ]
    y = 1.58
    for title, body in practices:
        rect(s, MARGIN, y + 0.05, 0.055, 0.2, fill=RED)
        text(s, MARGIN + 0.22, y, 7.3, 0.24, title, size=12.5, bold=True,
             color=INK)
        text(s, MARGIN + 0.22, y + 0.28, 7.25, 0.56, body, size=10.5,
             color=MUTED, line_spacing=1.24)
        y += 1.05

    text(s, 8.95, 2.35, 4.0, 0.7, "THANK YOU", size=38, bold=True, color=WHITE)
    accent_rule(s, 8.95, 3.16, w=1.5, color=WHITE, thickness=0.03)
    text(s, 8.95, 3.42, 3.95, 0.66,
         "Happy to go deeper on the architecture, the k6 numbers, or the firewall finding.",
         size=12, color=RGBColor(0xF2, 0xD6, 0xDB), line_spacing=1.3)

    text(s, 8.95, 4.60, 3.9, 0.22, "TEAM", size=9, bold=True,
         color=RGBColor(0xD8, 0x92, 0xA0), spacing=1.4)
    text(s, 8.95, 4.84, 3.9, 0.44, TEAM, size=25, bold=True, color=WHITE)
    text(s, 8.95, 5.36, 3.9, 0.24, MEMBERS, size=12, color=WHITE)
    text(s, 8.95, 5.78, 3.9, 0.24, "room107_Agileish  ·  branch dev", size=10,
         color=RGBColor(0xD8, 0x92, 0xA0))

    if asset("hsbc.png"):
        s.shapes.add_picture(asset("hsbc.png"), Inches(8.95), Inches(6.68),
                             height=Inches(0.4))
    else:
        text(s, 8.95, 6.72, 2.0, 0.3, "HSBC", size=13, bold=True, color=WHITE)

    hairline(s, MARGIN, 6.99, 7.52)
    text(s, MARGIN, 7.08, 6.0, 0.24, f"{TEAM}  ·  {PROJECT}", size=9,
         color=MUTED)
    text(s, 6.52, 7.08, 1.6, 0.24, f"{TOTAL_SLIDES} / {TOTAL_SLIDES}", size=9,
         color=MUTED, align=PP_ALIGN.RIGHT)
    return s


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    prs.core_properties.title = PROJECT
    prs.core_properties.author = f"Team {TEAM}"
    prs.core_properties.subject = "Client review — August 2026"

    slide_title(prs, blank)
    slide_problem(prs, blank)
    slide_architecture(prs, blank)
    slide_performance(prs, blank)
    slide_future(prs, blank)
    slide_close(prs, blank)

    prs.save(OUT)
    print(f"Wrote {OUT} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
